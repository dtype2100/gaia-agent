"""GAIA 첨부 파일 처리. IBM Docling, EasyOCR 및 고급 폴백 탑재.

개선점:
- 파일 다운로드 후 로컬 캐시 디렉토리 `.cache/attachments/{task_id}/{filename}` 에 저장.
- 도구 출력 시작 부분에 로컬 **절대 경로(Absolute Path)**를 포함하여 반환하여, 
  에이전트가 `exec_python_code`에서 Pandas, SQLite3, 또는 외부 스크립트로 파일을 다룰 수 있도록 지원.
- 문서 파싱:
  * 1차 시도: **IBM Docling**을 사용해 PDF, DOCX, PPTX, XLSX, 이미지를 레이아웃 보존 마크다운으로 변환.
  * 2차 폴백: Docling 실패 시 개별 라이브러리(python-docx, python-pptx, pypdf, pandas)로 텍스트/표 파싱.
  * 이미지 OCR: Docling 실패 또는 미지원 이미지의 경우 **EasyOCR**을 이용해 로컬 고정밀 텍스트 추출.
  * ZIP 아카이브: ZIP 파일 자동 압축 해제, 파일 목록 시각화, 내부 텍스트 파일(CSV, TXT 등) 본문 미리보기 제공.
"""
from __future__ import annotations

import io
import os
import re
import zipfile
import requests
from PIL import Image
from langchain_core.tools import tool

_DEFAULT_API_URL = "https://agents-course-unit4-scoring.hf.space"

_CURRENT_TASK = {"id": None, "question": None}
_QUESTION_INDEX: dict = {}


def prefetch_question_index() -> dict:
    """채점 서버 /questions 1회 호출 → {질문본문: task_id} 사전."""
    try:
        r = requests.get(f"{_DEFAULT_API_URL}/questions", timeout=15)
        r.raise_for_status()
        idx = {}
        for item in r.json():
            qt = (item.get("question") or "").strip()
            tid = item.get("task_id")
            if qt and tid:
                if qt in idx and idx[qt] != tid:
                    print(
                        "Warning: duplicate question text in prefetch index — "
                        f"task_id {idx[qt]!r} will be overwritten by {tid!r}"
                    )
                idx[qt] = tid
        return idx
    except Exception as e:
        print(f"Warning: could not prefetch question index: {e}")
        return {}


def set_question_index(idx: dict) -> None:
    global _QUESTION_INDEX
    _QUESTION_INDEX = idx


def set_current_task(question: str):
    tid = _QUESTION_INDEX.get(question.strip())
    _CURRENT_TASK["id"] = tid
    _CURRENT_TASK["question"] = question
    return tid


def _extract_filename(headers, url: str) -> str:
    cd = headers.get("Content-Disposition", "")
    m = re.search(r'filename\*?=(?:UTF-8\'\')?"?([^";\r\n]+)"?', cd)
    if m:
        return m.group(1).strip().strip('"')
    return url.rsplit("/", 1)[-1]


def _is_excel(content_type: str, ext: str) -> bool:
    if ext in ("xlsx", "xls"):
        return True
    ct = content_type.lower()
    return "spreadsheet" in ct or ct.endswith("xlsx") or ct.endswith("xls") or "excel" in ct


def _is_pdf(content_type: str, ext: str) -> bool:
    return ext == "pdf" or "pdf" in content_type.lower()


def _is_image(content_type: str, ext: str) -> bool:
    return ext in ("png", "jpg", "jpeg", "webp", "gif", "bmp") or content_type.lower().startswith("image/")


def _is_audio(content_type: str, ext: str) -> bool:
    return ext in ("mp3", "wav", "m4a", "ogg", "flac") or content_type.lower().startswith("audio/")


def _is_zip(content_type: str, ext: str) -> bool:
    return ext == "zip" or "zip" in content_type.lower() or "octet-stream" in content_type.lower() and ext == "zip"


# --- 1차 파서: IBM Docling ---
def _handle_docling(filepath: str) -> str:
    """IBM Docling을 이용해 문서 파일을 완벽한 마크다운 문서로 변환."""
    try:
        from docling.document_converter import DocumentConverter
        print(f"[Docling] Converting file: {filepath}")
        converter = DocumentConverter()
        result = converter.convert(filepath)
        md = result.document.export_to_markdown()
        return md
    except Exception as e:
        print(f"[Docling] Docling parsing failed: {e}. Falling back to default parsers.")
        return ""


# --- 2차 폴백 파서들 ---

def _handle_docx_fallback(filepath: str) -> str:
    """Word 파일 폴백 파서."""
    try:
        import docx
        doc = docx.Document(filepath)
        text_parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                text_parts.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                text_parts.append("| " + " | ".join(row_text) + " |")
        return "\n".join(text_parts)
    except Exception as e:
        return f"[DOCX Fallback Error] {e}"


def _handle_pptx_fallback(filepath: str) -> str:
    """PowerPoint 파일 폴백 파서."""
    try:
        import pptx
        prs = pptx.Presentation(filepath)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
        return "\n".join(text_parts)
    except Exception as e:
        return f"[PPTX Fallback Error] {e}"


def _handle_excel_fallback(filepath: str) -> str:
    """Excel 파일 폴백 파서."""
    try:
        import pandas as pd
        sheets = pd.read_excel(filepath, sheet_name=None)
        parts = []
        for name, df in sheets.items():
            parts.append(f"--- Sheet: {name} ---\n{df.to_csv(index=False)}")
        return "\n\n".join(parts)
    except Exception as e:
        return f"[Excel Fallback Error] {e}"


def _handle_pdf_fallback(filepath: str) -> str:
    """PDF 파일 폴백 파서."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(filepath)
        parts = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            parts.append(f"--- Page {i+1} ---\n{txt}")
        return "\n\n".join(parts)
    except Exception as e:
        return f"[PDF Fallback Error] {e}"


def _handle_zip(filepath: str, task_id: str) -> str:
    """ZIP 압축파일 자동 해제 및 텍스트 파일 미리보기 제공."""
    try:
        extract_dir = os.path.join(os.getcwd(), ".cache", "attachments", str(task_id), "extracted")
        os.makedirs(extract_dir, exist_ok=True)
        with zipfile.ZipFile(filepath, 'r') as zip_ref:
            zip_ref.extractall(extract_dir)

        file_list = []
        text_contents = []
        for root, dirs, files in os.walk(extract_dir):
            for file in files:
                full_path = os.path.join(root, file)
                rel_path = os.path.relpath(full_path, extract_dir)
                file_list.append(rel_path)
                
                # 텍스트 계열 파일은 2000글자까지 미리보기 수집
                ext = file.rsplit(".", 1)[-1].lower() if "." in file else ""
                if ext in ("txt", "csv", "json", "py", "html", "css", "md", "xml", "yaml", "yml", "ini", "log"):
                    try:
                        with open(full_path, "r", encoding="utf-8", errors="ignore") as f:
                            content = f.read(2000)
                            if len(content) >= 2000:
                                content += "\n...[truncated]"
                            text_contents.append(f"=== File: {rel_path} (Absolute Path: {os.path.abspath(full_path)}) ===\n{content}")
                    except Exception as fe:
                        text_contents.append(f"=== File: {rel_path} ===\n(Read Error: {fe})")

        summary = f"ZIP Archive extracted successfully.\nExtracted to absolute folder: {os.path.abspath(extract_dir)}\n\nFiles list:\n"
        summary += "\n".join([f"- {f}" for f in file_list])
        if text_contents:
            summary += "\n\n--- Nested Text-Based Files Content (Previews) ---\n\n" + "\n\n".join(text_contents)
        return summary
    except Exception as e:
        return f"[ZIP Extraction Error] {e}"


def _handle_image_ocr(filepath: str) -> str:
    """EasyOCR을 통한 이미지 텍스트 고정밀 추출."""
    try:
        import easyocr
        import torch
        print(f"[EasyOCR] Extracting text from image: {filepath}")
        
        # ZeroGPU/GPU 환경 여부에 따라 gpu 활성화
        reader = easyocr.Reader(['en'], gpu=torch.cuda.is_available())
        results = reader.readtext(filepath, detail=0)
        ocr_text = "\n".join(results)
        return f"### OCR Extracted Text:\n{ocr_text}"
    except Exception as e:
        return f"[Image OCR Fallback Error] {e}"


def _handle_image(filepath: str, content_type: str) -> str:
    """이미지 처리 메인 진입점. 1차 Docling -> 2차 EasyOCR."""
    meta = ""
    try:
        img = Image.open(filepath)
        w, h = img.size
        meta = f"[Image metadata: {w}x{h}px, {content_type}]\n"
    except Exception:
        pass

    # Docling 시도
    docling_result = _handle_docling(filepath)
    if docling_result.strip():
        return meta + docling_result

    # EasyOCR 시도
    ocr_result = _handle_image_ocr(filepath)
    return meta + ocr_result


def _handle_audio(filepath: str, content_type: str) -> str:
    """V0: 오디오 전사 미구현 안내 및 메타데이터 반환."""
    return (
        f"[Audio attached (Content-Type: {content_type}, File saved locally). "
        f"NOTE: audio transcription is not yet wired in this runtime. If the "
        f"question REQUIRES transcribing this audio, you can write a python script "
        f"utilizing the local absolute path below, or try to decode it yourself.]"
    )


def download_attachment(task_id: str) -> dict | None:
    """task_id에 해당하는 첨부파일을 다운로드하여 로컬에 캐시하고,
    파일 메타데이터(절대경로, filename, content_type, size, ext)를 반환.
    파일이 없거나 404인 경우 None 반환.
    """
    try:
        task_dir = os.path.join(os.getcwd(), ".cache", "attachments", str(task_id))
        if os.path.exists(task_dir):
            files = [f for f in os.listdir(task_dir) if os.path.isfile(os.path.join(task_dir, f))]
            if files:
                filename = files[0]
                filepath = os.path.join(task_dir, filename)
                abs_filepath = os.path.abspath(filepath)
                ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
                
                content_type = ""
                if ext in ("png", "jpg", "jpeg", "webp", "gif", "bmp"):
                    content_type = f"image/{ext}"
                elif ext in ("mp3", "wav", "m4a", "ogg", "flac"):
                    content_type = f"audio/{ext}"
                elif ext == "pdf":
                    content_type = "application/pdf"
                elif ext == "zip":
                    content_type = "application/zip"
                
                return {
                    "abs_path": abs_filepath,
                    "filename": filename,
                    "content_type": content_type,
                    "size": os.path.getsize(abs_filepath),
                    "ext": ext
                }

        url = f"{_DEFAULT_API_URL}/files/{task_id}"
        print(f"[Attachments] Downloading attached file from: {url}")
        r = requests.get(url, timeout=30)
        if r.status_code == 404:
            return None
        r.raise_for_status()

        content_type = r.headers.get("Content-Type", "")
        filename = _extract_filename(r.headers, url)
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        os.makedirs(task_dir, exist_ok=True)
        filepath = os.path.join(task_dir, filename)
        with open(filepath, "wb") as f:
            f.write(r.content)
        abs_filepath = os.path.abspath(filepath)

        return {
            "abs_path": abs_filepath,
            "filename": filename,
            "content_type": content_type,
            "size": len(r.content),
            "ext": ext
        }
    except Exception as e:
        print(f"Error in download_attachment for task_id {task_id}: {e}")
        return None


@tool
def get_attached_file() -> str:
    """Download the file attached to the CURRENT GAIA task, save it locally, and return its content.
    Takes no arguments — the current task_id is auto-resolved from the question context.

    Use this whenever the question references a file, spreadsheet, image, audio, PDF, code listing,
    CSV, or any external resource. Returns:
      - The local absolute path to the downloaded file (allowing python execution tools to open it directly).
      - Markdown representation of documents (via IBM Docling) and tables.
      - Extracted OCR text from images (via EasyOCR).
      - Directory listing and text previews for ZIP archives.
    """
    task_id = _CURRENT_TASK.get("id")
    if not task_id:
        return "No task context available — likely no file attached for this question."
    try:
        meta = download_attachment(task_id)
        if not meta:
            return "No file attached to this task."

        abs_filepath = meta["abs_path"]
        filename = meta["filename"]
        content_type = meta["content_type"]
        ext = meta["ext"]

        header_info = f"### [Local Save Status] File saved locally successfully.\n" \
                      f"- **Local Absolute Path**: `{abs_filepath}`\n" \
                      f"- **Filename**: `{filename}`\n" \
                      f"- **Content-Type**: `{content_type}`\n" \
                      f"- **File Size**: `{meta['size']}` bytes\n\n"

        # --- 파싱 로직 분기 ---

        # 1. ZIP 아카이브
        if _is_zip(content_type, ext):
            zip_content = _handle_zip(abs_filepath, task_id)
            return header_info + zip_content

        # 2. 오디오 파일
        if _is_audio(content_type, ext):
            audio_content = _handle_audio(abs_filepath, content_type)
            return header_info + audio_content

        # 3. 이미지 파일
        if _is_image(content_type, ext):
            image_content = _handle_image(abs_filepath, content_type)
            return header_info + image_content

        # 4. 문서 계열 (PDF, Word, PPTX, Excel)
        if _is_pdf(content_type, ext) or _is_excel(content_type, ext) or ext in ("docx", "pptx"):
            # 1차 Docling 변환 시도
            docling_text = _handle_docling(abs_filepath)
            if docling_text.strip():
                if len(docling_text) > 15000:
                    docling_text = docling_text[:15000] + "\n...[truncated]"
                return header_info + docling_text

            # 실패 시 개별 폴백 작동
            print(f"[Parser Fallback] Docling returned empty, falling back...")
            if _is_pdf(content_type, ext):
                parsed_text = _handle_pdf_fallback(abs_filepath)
            elif _is_excel(content_type, ext):
                parsed_text = _handle_excel_fallback(abs_filepath)
            elif ext == "docx":
                parsed_text = _handle_docx_fallback(abs_filepath)
            elif ext == "pptx":
                parsed_text = _handle_pptx_fallback(abs_filepath)
            else:
                parsed_text = "(Fallback parser not found)"

            if len(parsed_text) > 12000:
                parsed_text = parsed_text[:12000] + "\n...[truncated]"
            return header_info + parsed_text

        # 5. 일반 텍스트 계열 파일 (CSV, JSON, Python, PDB 등)
        try:
            with open(abs_filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            if len(text) > 12000:
                text = text[:12000] + "\n...[truncated]"
            return header_info + f"### File Content:\n```\n{text}\n```"
        except Exception:
            pass

        return header_info + f"Binary file (cannot display as plain text directly). " \
                             f"Use python environment to inspect: `open('{abs_filepath.replace(chr(92), '/')}', 'rb')`"

    except Exception as e:
        return f"get_attached_file error: {e}"
